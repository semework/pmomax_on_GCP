#!/usr/bin/env python3
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from docx import Document
except Exception as e:
    print('Missing dependencies:', e)
    print('Install with: pip install python-pptx python-docx')
    sys.exit(1)

root = Path(__file__).resolve().parent.parent
md_path = root / 'docs' / 'PMOMax_Business_Pricing.md'
if not md_path.exists():
    print('Markdown source not found:', md_path)
    sys.exit(1)

text = md_path.read_text(encoding='utf-8')

# Simple PPTX: create slides for main headings and summary tables
prs = Presentation()

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'PMOMax — Business & Pricing Overview'
subtitle.text = 'Generated from PMOMax_Business_Pricing.md'

# Add a few summary slides: break on H2 ("## ")
sections = []
current = None
for line in text.splitlines():
    if line.startswith('## '):
        if current:
            sections.append(current)
        current = {'title': line[3:].strip(), 'body': []}
    elif line.startswith('# '):
        continue
    else:
        if current is None:
            current = {'title': 'Overview', 'body': []}
        current['body'].append(line)
if current:
    sections.append(current)

for sec in sections[:8]:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = sec['title']
    body = slide.shapes.placeholders[1].text_frame
    # take first 6 non-empty lines
    lines = [l.strip() for l in sec['body'] if l.strip()]
    for i, l in enumerate(lines[:8]):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = l
        p.level = 0

out_ppt = root / 'docs' / 'PMOMax_Business_Pricing.pptx'
prs.save(out_ppt)
print('Wrote PPTX:', out_ppt)

# Create a Word document with the same markdown split into sections
from docx.shared import Pt

doc = Document()
doc.add_heading('PMOMax — Business & Pricing Overview', level=1)
for line in text.splitlines():
    if line.startswith('# '):
        continue
    elif line.startswith('## '):
        doc.add_heading(line[3:].strip(), level=2)
    elif line.startswith('- '):
        doc.add_paragraph(line[2:].strip(), style='List Bullet')
    elif line.strip() == '':
        doc.add_paragraph('')
    else:
        doc.add_paragraph(line)

out_docx = root / 'docs' / 'PMOMax_Business_Pricing.docx'
doc.save(out_docx)
print('Wrote DOCX:', out_docx)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Flowchart nodes and edges (simplified for PowerPoint)
nodes = [
    ("Customer Tenant (User)", 3.5, 0.5, RGBColor(0xE3, 0xF2, 0xFD)),
    ("Frontend (React SPA)", 3.5, 2.0, RGBColor(0xF5, 0xF5, 0xF5)),
    ("Partner Tenant (Backend API)", 3.5, 3.5, RGBColor(0xFF, 0xF9, 0xC4)),
    ("AI/Parsing Logic (Gemini, Local)", 3.5, 5.0, RGBColor(0xF3, 0xE5, 0xF5)),
    ("Google Cloud Platform (Cloud Run / GKE)", 3.5, 6.5, RGBColor(0xE8, 0xF5, 0xE9)),
]
edges = [
    (0, 1, "Uploads Text/Document"),
    (1, 2, "Sends Data"),
    (2, 3, "Parse/AI/Export"),
    (3, 2, "Returns PID/Reply"),
    (2, 1, "Returns Data"),
    (1, 0, "Shows PID/Assistant"),
    (1, 2, "Export Request"),
    (2, 1, "Export File"),
    (2, 4, "Compliance/Health"),
    (4, 1, "Hosts"),
    (4, 2, "Hosts"),
]

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Draw nodes
shapes = []
for i, (label, x, y, color) in enumerate(nodes):
    left = Inches(x)
    top = Inches(y)
    width = Inches(2.2)
    height = Inches(0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = label
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shapes.append(shape)



# Draw arrows (lines with arrowheads and labels, vertical layout)
for src, dst, label in edges:
    src_shape = shapes[src]
    dst_shape = shapes[dst]
    x1 = src_shape.left + src_shape.width / 2
    y1 = src_shape.top + src_shape.height
    x2 = dst_shape.left + dst_shape.width / 2
    y2 = dst_shape.top
    # Draw a line (arrow)
    line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x1, y1, 0, y2-y1)
    line.line.width = Pt(2)
    line.line.end_arrowhead.style = 2  # Arrowhead
    # Add label as a text box near the middle of the arrow
    mid_x = x1 - Inches(1)
    mid_y = (y1 + y2) / 2 - Inches(0.2)
    textbox = slide.shapes.add_textbox(mid_x, mid_y, Inches(2), Inches(0.3))
    textbox.text = label
    textbox.text_frame.paragraphs[0].font.size = Pt(10)
    textbox.text_frame.paragraphs[0].font.name = "Arial"

prs.save("/Users/mulugetasemework/Library/Mobile Documents/com~apple~CloudDocs/projects/PMO/pmo01/PMOMax_App_Flowchart.pptx")
